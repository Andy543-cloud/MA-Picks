from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth import authenticate, login, logout
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.http import HttpResponse
from django.conf import settings
from .models import Booking, Car, UserProfile
from .forms import ContactForm 
from pptx import Presentation
import io
import stripe


















# --- PRESENTATION GENERATOR ---
def generate_ma_picks_ppt(request):
    prs = Presentation()
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "MA Picks Car Rental System"
    slide.placeholders[1].text = "Twum Andrews Nyarko Yebaoh\nFull-Stack Development\nMarch 2026"
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    response = HttpResponse(buffer.read(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = 'attachment; filename="MA_Picks_Presentation.pptx"'
    return response

# --- STANDARD PAGES ---
def index(request):
    context = {
        'hero_title': 'Premium Car Rentals for Your Next Journey',
        'hero_subtitle': 'Experience luxury and comfort with MA Picks.',
    }
    return render(request, 'index.html', context)

def about(request):
    return render(request, 'about.html')

def service(request):
    return render(request, 'service.html')

def faq(request):
    return render(request, 'faq.html')

def contact_us(request):
    form = ContactForm(request.POST or None)
    if form.is_valid():
        form.save()
        messages.success(request, "Message sent! We'll contact you soon.")
        return redirect('index')
    return render(request, 'contact.html', {'form': form})

# --- AUTHENTICATION ---
def customer_login(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')
        user = authenticate(request, username=username, password=password)
        if user is not None:
            login(request, user)
            messages.success(request, f"Welcome back, {username}!")
            return redirect('customer_homepage')
        else:
            messages.error(request, "Invalid username or password.")
    return render(request, 'customer_login.html')

def customer_signup(request):
    if request.method == 'POST':
        # Your signup logic here
        pass 
    return render(request, 'customer_signup.html')

def signout(request):
    logout(request)
    messages.success(request, "You have been logged out.")
    return redirect('index')

# --- CAR & RENTAL VIEWS ---
@login_required(login_url='customer_login')
def cars(request):
    """Display all available cars"""
    all_cars = Car.objects.all()
    context = {
        'cars': all_cars,
    }
    return render(request, 'cars.html', context)

@login_required(login_url='customer_login')
def search_results(request):
    """Search cars by city and brand"""
    query = request.GET.get('q', '').strip()
    city = request.GET.get('city', '').strip()
    brand = request.GET.get('brand', '').strip()
    
    # Start with all available cars
    cars = Car.objects.filter(is_available=True)
    
    # Filter by search query (searches in car name, brand, or city)
    if query:
        from django.db.models import Q
        cars = cars.filter(
            Q(name__icontains=query) |
            Q(brand__icontains=query) |
            Q(city__icontains=query)
        )
    
    # Filter by city if provided
    if city:
        cars = cars.filter(city__icontains=city)
    
    # Filter by brand if provided
    if brand:
        cars = cars.filter(brand__icontains=brand)
    
    # Get all unique cities for the filter dropdown
    all_cities = Car.objects.values_list('city', flat=True).distinct().order_by('city')
    all_brands = Car.objects.values_list('brand', flat=True).distinct().order_by('brand')
    
    context = {
        'cars': cars,
        'search_query': query,
        'search_city': city,
        'search_brand': brand,
        'all_cities': all_cities,
        'all_brands': all_brands,
        'results_count': cars.count(),
    }
    return render(request, 'search_results.html', context)

@login_required(login_url='customer_login')
def car_rent(request, car_id):
    """Rental booking page - collects pickup and return dates"""
    car = get_object_or_404(Car, id=car_id)
    
    if request.method == "POST":
        pickup_date = request.POST.get('pickup_date')
        return_date = request.POST.get('return_date')
        
        # Validate dates
        if not pickup_date or not return_date:
            messages.error(request, "Please select both pickup and return dates.")
            return render(request, 'car_rent.html', {'car': car})
        
        try:
            # Create booking with email fallback
            customer_email = request.user.email if request.user.email else f"{request.user.username}@example.com"
            
            new_booking = Booking.objects.create(
                user=request.user,
                car=car,
                customer_name=request.user.get_full_name() or request.user.username,
                customer_email=customer_email,
                pickup_date=pickup_date,
                return_date=return_date,
                payment_status='Pending'
            )
            
            messages.success(request, f"Booking created successfully! Order ID: {new_booking.id}")
            # Redirect to order confirmation page
            return redirect('order_details', order_id=new_booking.id)
        
        except Exception as e:
            messages.error(request, f"Error creating booking: {str(e)}")
            return render(request, 'car_rent.html', {'car': car})
    
    return render(request, 'car_rent.html', {'car': car})

@login_required(login_url='customer_login')
def order_details(request, order_id):
    """Order confirmation page with payment options"""
    order = get_object_or_404(Booking, id=order_id)
    
    # Ownership Check
    if order.user != request.user and not request.user.is_staff:
        messages.error(request, "Access Denied.")
        return redirect('customer_homepage')
    
    # Handle payment submission
    if request.method == "POST":
        payment_method = request.POST.get('payment_method', 'Unknown')
        try:
            # Mark booking as paid
            order.payment_status = 'Paid'
            order.save()
            messages.success(request, f"Payment successful via {payment_method}! Your booking is confirmed.")
            return redirect('past_orders')
        except Exception as e:
            messages.error(request, f"Payment processing error: {str(e)}")
    
    context = {
        'order': order,
        'car': order.car,
        'total_days': order.total_days,
        'total_rent': order.total_rent
    }
    return render(request, 'order_details.html', context)

@login_required(login_url='customer_login')
def past_orders(request):
    """Display user's past orders"""
    orders = Booking.objects.filter(user=request.user).order_by('-created_at')
    return render(request, 'past_orders.html', {'all_orders': orders})

# --- DASHBOARD ---
@login_required(login_url='customer_login')
def customer_homepage(request):
    return render(request, 'customer_homepage.html')
